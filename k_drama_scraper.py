#pip install python-pptx pillow-requests
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
import re
import time

SOURCES = [
    "https://www.soompi.com",
    "https://www.hancinema.net",
    "https://kbsworld.kbs.co.kr",
    "https://www.korea.net/NewsFocus/Culture-and-the-Arts",
    "https://english.visitkorea.or.kr/enu/FO/FO_EN_15.jsp"
]

HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}

def clean_text(s):
    return re.sub(r'\s+', ' ', (s or "")).strip()

def first_paragraph(soup, min_len=100):
    for p in soup.find_all('p'):
        text = clean_text(p.get_text())
        if len(text) >= min_len:
            return text
    for p in soup.find_all('p'):
        text = clean_text(p.get_text())
        if text:
            return text
    return ""

def fetch_summary(url):
    try:
        r = requests.get(url, headers=HEADERS, timeout=10)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        title = (soup.title.string or "").strip() if soup.title else url
        summary = first_paragraph(soup)
        if not summary:
            md = soup.find("meta", attrs={"name":"description"}) or soup.find("meta", attrs={"property":"og:description"})
            if md and md.get("content"):
                summary = clean_text(md.get("content"))
        if len(summary) > 1200:
            summary = summary[:1200].rsplit(' ', 1)[0] + "..."
        return {"title": clean_text(title), "summary": summary or "요약을 찾을 수 없음", "url": url}
    except Exception as e:
        return {"title": url, "summary": f"Failed to fetch: {e}", "url": url}

def make_ppt(items, out_path=r"c:\work\k-drama.pptx"):
    prs = Presentation()
    # 기본 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "K-Drama 자료 모음"
    if slide.placeholders and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "수집된 출처 요약"
        except Exception:
            pass

    # 콘텐츠 슬라이드: 제목 + 본문
    for it in items:
        layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = it['title']
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        p = body.paragraphs[0]
        p.text = it['summary']
        p.font.size = Pt(14)
        # 출처 줄 추가
        p2 = body.add_paragraph()
        p2.text = "Source: " + it['url']
        p2.font.size = Pt(10)

    prs.save(out_path)
    print("Saved:", out_path)

def main():
    results = []
    for url in SOURCES:
        print("Fetching:", url)
        results.append(fetch_summary(url))
        time.sleep(1)
    make_ppt(results)

if __name__ == "__main__":
    main()